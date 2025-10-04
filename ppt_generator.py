# ppt_generator.py

import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import requests
import io
import re

def create_presentation(slides_data):
    """Builds a PowerPoint presentation from structured slide data."""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = slides_data[0]['title']
    subtitle.text = slides_data[0]['subtitle']

    # Add content slides
    for slide_info in slides_data[1:]:
        content_slide_layout = prs.slide_layouts[5] # Blank layout
        slide = prs.slides.add_slide(content_slide_layout)

        # Add Title
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.5))
        p = title_shape.text_frame.add_paragraph()
        p.text = slide_info['title']
        p.font.bold = True
        p.font.size = 320000

        # Add Content (Bullet Points)
        content_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(7), Inches(6))
        content_frame = content_shape.text_frame
        content_frame.word_wrap = True
        for point in slide_info['content']:
            clean_point = re.sub(r'^[*-]\s*', '', point).strip()
            if clean_point:
                p = content_frame.add_paragraph()
                p.text = clean_point
                p.level = 0
                p.font.size = 180000

        # Add Image
        if slide_info.get('image_url'):
            try:
                response = requests.get(slide_info['image_url'])
                response.raise_for_status()
                image_stream = io.BytesIO(response.content)
                slide.shapes.add_picture(image_stream, Inches(8.5), Inches(2), width=Inches(7))
            except requests.exceptions.RequestException as e:
                st.warning(f"Failed to download image for slide '{slide_info['title']}': {e}")

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio